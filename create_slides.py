"""
create_slides.py

Generate a PowerPoint presentation for:
  "Adaptive CI Heuristics for Route Optimization Under Near Real-Time Constraints"

Run:
    python create_slides.py
Output:
    presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# -----------------------------------------------------------------------
# Color palette
# -----------------------------------------------------------------------
NAVY      = RGBColor(0x1B, 0x27, 0x5A)   # dark navy — headings
BLUE      = RGBColor(0x26, 0x5E, 0xAD)   # medium blue — accents
LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xFF)   # near-white blue — bg tint
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x1A, 0x1A, 0x1A)
GRAY      = RGBColor(0x55, 0x55, 0x55)
GREEN     = RGBColor(0x1A, 0x7A, 0x3C)
ORANGE    = RGBColor(0xD4, 0x6A, 0x00)
RED       = RGBColor(0xB8, 0x1C, 0x1C)

# -----------------------------------------------------------------------
# Helpers
# -----------------------------------------------------------------------

def add_slide(prs, layout_index=6):
    layout = prs.slide_layouts[layout_index]
    return prs.slides.add_slide(layout)


def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=BLACK,
                align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_title_area(slide, title, subtitle=None):
    """Add a navy header band with title."""
    add_rect(slide, 0, 0, 10, 1.15, NAVY)
    add_textbox(slide, title, 0.3, 0.18, 9.4, 0.85,
                font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle, 0.3, 0.72, 9.4, 0.4,
                    font_size=13, color=RGBColor(0xBB, 0xCC, 0xFF), align=PP_ALIGN.LEFT)


def add_bullet_box(slide, title, bullets, left, top, width, height,
                   title_color=NAVY, bullet_color=BLACK, title_size=15, bullet_size=13):
    """Titled bullet-point box."""
    add_rect(slide, left, top, width, 0.38, NAVY)
    add_textbox(slide, title, left + 0.1, top + 0.04, width - 0.2, 0.3,
                font_size=title_size, bold=True, color=WHITE)

    txBox = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.42),
        Inches(width - 0.2), Inches(height - 0.55)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = b
        run.font.size = Pt(bullet_size)
        run.font.color.rgb = bullet_color


def add_metric_box(slide, label, value, sublabel, left, top, color=BLUE):
    add_rect(slide, left, top, 2.1, 1.35, color)
    add_textbox(slide, value, left, top + 0.12, 2.1, 0.65,
                font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, left, top + 0.72, 2.1, 0.35,
                font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, sublabel, left, top + 1.02, 2.1, 0.28,
                font_size=10, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)


# -----------------------------------------------------------------------
# Slides  (storytelling rewrite)
# -----------------------------------------------------------------------

def slide_title(prs):
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_rect(slide, 0, 0, 10, 7.5, NAVY)

    # Accent bar
    add_rect(slide, 0, 5.3, 10, 0.06, BLUE)
    # Tag line bar at very top
    add_rect(slide, 0, 0, 10, 0.5, RGBColor(0x12, 0x1A, 0x40))
    add_textbox(slide, "CSc 8810  ·  Computational Intelligence  ·  Georgia State University  ·  Spring 2026",
                0.3, 0.06, 9.4, 0.35, font_size=11,
                color=RGBColor(0x77, 0x88, 0xBB), align=PP_ALIGN.CENTER)

    add_textbox(slide,
        "FACI-DVRP",
        0.6, 0.7, 8.8, 1.0, font_size=44, bold=True,
        color=RGBColor(0x66, 0xAA, 0xFF), align=PP_ALIGN.CENTER)

    add_textbox(slide,
        "Fuzzy-Adaptive Computational Intelligence\nfor Dynamic Vehicle Routing",
        0.6, 1.65, 8.8, 1.1, font_size=22, bold=False,
        color=WHITE, align=PP_ALIGN.CENTER)

    # Divider
    add_rect(slide, 2.5, 2.95, 5.0, 0.04, RGBColor(0x44, 0x66, 0xAA))

    add_textbox(slide,
        "3 CI techniques  ·  Real-time routing  ·  Solomon benchmark",
        0.6, 3.1, 8.8, 0.45, font_size=13,
        color=RGBColor(0x88, 0xAA, 0xDD), align=PP_ALIGN.CENTER)

    add_textbox(slide,
        "Pranay Kukkadapu",
        0.6, 3.75, 8.8, 0.42, font_size=16, bold=True,
        color=RGBColor(0xBB, 0xCC, 0xFF), align=PP_ALIGN.CENTER)

    add_textbox(slide,
        "Masters in Data Science and Analytics  ·  Georgia State University",
        0.6, 4.18, 8.8, 0.38, font_size=12,
        color=RGBColor(0x88, 0x99, 0xCC), align=PP_ALIGN.CENTER)


def slide_hook(prs):
    """Story hook — the delivery company scenario."""
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_area(slide, "Picture This...", "A day in the life of a delivery company")

    # Big scenario text
    add_textbox(slide,
        "It's 8 AM. Your dispatcher has planned perfect routes for 10 drivers.",
        0.4, 1.3, 9.2, 0.55, font_size=16, bold=True, color=NAVY)

    # Three disruption cards
    disruptions = [
        ("9:15 AM",   "New customer\ncalls in",
         "A last-minute order arrives after trucks have already left.",
         RGBColor(0x26, 0x5E, 0xAD)),
        ("10:30 AM",  "Traffic jam\non Route 9",
         "An accident blocks the road that 3 drivers are scheduled to use.",
         ORANGE),
        ("11:45 AM",  "Heavy rain\nbegins",
         "Speeds drop 30% across the entire city. Every route is now slower.",
         RGBColor(0x1A, 0x6A, 0x8A)),
    ]
    for i, (time, title, desc, color) in enumerate(disruptions):
        lx = 0.3 + i * 3.18
        ty = 2.05
        add_rect(slide, lx, ty, 3.0, 0.35, color)
        add_textbox(slide, time, lx + 0.1, ty + 0.05, 2.8, 0.25,
                    font_size=11, bold=True, color=WHITE)
        add_rect(slide, lx, ty + 0.35, 3.0, 2.6, LIGHT_BG, color)
        add_textbox(slide, title, lx + 0.12, ty + 0.48, 2.76, 0.65,
                    font_size=15, bold=True, color=color)
        add_textbox(slide, desc, lx + 0.12, ty + 1.18, 2.76, 1.6,
                    font_size=11, color=BLACK)

    # The question
    add_rect(slide, 0.3, 5.1, 9.4, 0.06, BLUE)
    add_textbox(slide,
        "What do you do?  You can't ignore it — but you also can't spend 10 minutes recalculating\n"
        "everything while the driver is waiting on the side of the road.",
        0.4, 5.28, 9.2, 0.8, font_size=13, color=NAVY, align=PP_ALIGN.CENTER)

    add_rect(slide, 2.0, 6.2, 6.0, 0.55, NAVY)
    add_textbox(slide, "This is the Dynamic Vehicle Routing Problem (DVRP)",
                2.0, 6.27, 6.0, 0.4, font_size=13, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)


def slide_two_bad_options(prs):
    """The naive options and why both fail."""
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_area(slide, "The Two Obvious Answers", "And why both of them fail")

    # Option 1 — Ignore it
    add_rect(slide, 0.3, 1.3, 4.4, 0.48, ORANGE)
    add_textbox(slide, 'Option A:  "Just ignore it"', 0.42, 1.37, 4.2, 0.34,
                font_size=15, bold=True, color=WHITE)
    add_rect(slide, 0.3, 1.78, 4.4, 3.0, RGBColor(0xFF, 0xF3, 0xE0), ORANGE)
    ignore_pts = [
        "Keep the original plan. Don't update anything.",
        "",
        "  ✗  New customers never get served",
        "  ✗  Traffic jams make routes take 2x longer",
        "  ✗  Costs spiral as disruptions pile up",
        "  ✗  0% acceptance rate in our experiment",
    ]
    add_textbox(slide, "\n".join(ignore_pts), 0.45, 1.9, 4.1, 2.7,
                font_size=12, color=BLACK)

    # VS
    add_textbox(slide, "VS", 4.75, 2.9, 0.5, 0.6, font_size=22, bold=True,
                color=NAVY, align=PP_ALIGN.CENTER)

    # Option 2 — Rebuild everything
    add_rect(slide, 5.3, 1.3, 4.4, 0.48, RED)
    add_textbox(slide, 'Option B:  "Rebuild everything"', 5.42, 1.37, 4.2, 0.34,
                font_size=15, bold=True, color=WHITE)
    add_rect(slide, 5.3, 1.78, 4.4, 3.0, RGBColor(0xFF, 0xEE, 0xEE), RED)
    rebuild_pts = [
        "Solve the entire problem from scratch every time.",
        "",
        "  ✗  Takes 100 – 400 ms per event (too slow)",
        "  ✗  Drivers get completely new routes mid-trip",
        "  ✗  Route stability collapses to 0.41",
        "  ✗  Impractical in real operations",
    ]
    add_textbox(slide, "\n".join(rebuild_pts), 5.45, 1.9, 4.1, 2.7,
                font_size=12, color=BLACK)

    # Punchline
    add_rect(slide, 0.3, 5.0, 9.4, 0.06, BLUE)
    add_rect(slide, 0.3, 5.15, 9.4, 0.92, NAVY)
    add_textbox(slide,
        "We need something smarter — fast enough for real-time, good enough to be useful.\n"
        "That's exactly what FACI-DVRP is built to do.",
        0.5, 5.22, 9.0, 0.8, font_size=13, color=WHITE, align=PP_ALIGN.CENTER)


def slide_solution_overview(prs):
    """Introduce FACI-DVRP as the answer with 3 CI techniques."""
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_area(slide, "Meet FACI-DVRP", "Three CI techniques working together")

    add_textbox(slide,
        "Fuzzy-Adaptive Computational Intelligence for Dynamic Vehicle Routing Problems",
        0.4, 1.22, 9.2, 0.42, font_size=13, color=GRAY, align=PP_ALIGN.CENTER)

    # Phase labels
    add_textbox(slide, "BEFORE THE DAY", 0.4, 1.78, 3.2, 0.32,
                font_size=11, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    add_textbox(slide, "WHEN DISRUPTION HITS", 3.9, 1.78, 2.4, 0.32,
                font_size=11, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    add_textbox(slide, "FIX IT (FAST)", 6.6, 1.78, 3.1, 0.32,
                font_size=11, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

    # CI Technique boxes
    boxes = [
        (0.3,  2.15, 3.3, BLUE,   "CI Technique 1",  "Ant Colony\nOptimization (ACO)",
         "Inspired by ants finding shortest paths.\nBuilds high-quality initial routes\nbefore the day starts."),
        (3.85, 2.15, 2.5, GREEN,  "CI Technique 2",  "Fuzzy Logic\nController",
         "Reasons like a human:\n'How severe is this?'\nAllocates repair time accordingly."),
        (6.55, 2.15, 3.15, ORANGE, "CI Technique 3",  "2-opt Local\nSearch + Repair",
         "Improves routes by trying\nsmart swaps — within the\ntime budget given by fuzzy."),
    ]

    for lx, ty, w, color, tag, title, desc in boxes:
        add_rect(slide, lx, ty, w, 0.32, RGBColor(0xDD, 0xEE, 0xFF))
        add_textbox(slide, tag, lx + 0.08, ty + 0.04, w - 0.16, 0.24,
                    font_size=10, bold=True, color=color)
        add_rect(slide, lx, ty + 0.32, w, 0.62, color)
        add_textbox(slide, title, lx + 0.1, ty + 0.38, w - 0.2, 0.52,
                    font_size=14, bold=True, color=WHITE)
        add_rect(slide, lx, ty + 0.94, w, 2.1, LIGHT_BG, color)
        add_textbox(slide, desc, lx + 0.12, ty + 1.06, w - 0.24, 1.85,
                    font_size=11, color=BLACK)

    # Arrows between boxes
    add_textbox(slide, "→", 3.33, 3.3, 0.5, 0.5, font_size=22, bold=True,
                color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(slide, "→", 6.1, 3.3, 0.44, 0.5, font_size=22, bold=True,
                color=NAVY, align=PP_ALIGN.CENTER)

    # Key principle
    add_rect(slide, 0.3, 5.45, 9.4, 0.08, BLUE)
    add_rect(slide, 0.3, 5.6, 9.4, 0.82, NAVY)
    add_textbox(slide,
        "Key principle: Never rebuild the entire plan. Find what's broken, fix only that — "
        "and only spend as much time as the disruption actually deserves.",
        0.5, 5.67, 9.0, 0.7, font_size=13, color=WHITE, align=PP_ALIGN.CENTER)


def slide_aco(prs):
    """ACO explained like a story."""
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_area(slide, "Technique 1: Ant Colony Optimization",
                   "How ant colonies inspire better route planning")

    # Analogy
    add_rect(slide, 0.3, 1.25, 9.4, 0.06, BLUE)
    add_textbox(slide,
        "Real ants find the shortest path to food by laying pheromone trails — "
        "the more ants use a path, the stronger the trail, and the more ants follow it.",
        0.4, 1.38, 9.2, 0.55, font_size=13, color=NAVY, align=PP_ALIGN.CENTER)

    # Steps
    steps = [
        ("Send out the ants",
         "10 virtual ants each build a complete\nrouting plan from scratch, choosing\ncustomers probabilistically based on\ndistance and pheromone strength."),
        ("Reward the good paths",
         "After each round, the best solution\ngets extra pheromone deposited on\nits edges — making those paths more\nlikely to be chosen next time."),
        ("Evaporate the bad ones",
         "All pheromone slowly fades (rate 0.5).\nPaths that aren't reinforced disappear.\nThe colony converges on the best\nroutes over 20 iterations."),
    ]
    for i, (title, desc) in enumerate(steps):
        lx = 0.3 + i * 3.22
        ty = 2.1
        add_rect(slide, lx, ty, 3.0, 0.46, BLUE)
        add_textbox(slide, f"Step {i+1}:  {title}", lx + 0.1, ty + 0.08, 2.8, 0.32,
                    font_size=12, bold=True, color=WHITE)
        add_rect(slide, lx, ty + 0.46, 3.0, 2.15, LIGHT_BG, BLUE)
        add_textbox(slide, desc, lx + 0.12, ty + 0.6, 2.76, 1.9,
                    font_size=11, color=BLACK)

    # Formula
    add_rect(slide, 0.3, 4.65, 5.5, 0.38, NAVY)
    add_textbox(slide, "Selection formula", 0.42, 4.7, 5.2, 0.28,
                font_size=11, bold=True, color=WHITE)
    add_rect(slide, 0.3, 5.03, 5.5, 0.62, LIGHT_BG, NAVY)
    add_textbox(slide, "P(next customer j)  ∝  τ[i,j]^α  ×  (1/distance)^β",
                0.45, 5.12, 5.2, 0.45, font_size=12, color=BLACK)

    # Result
    add_rect(slide, 6.0, 4.65, 3.7, 1.0, NAVY)
    add_textbox(slide, "Result on C101 benchmark:", 6.12, 4.72, 3.5, 0.3,
                font_size=11, bold=True, color=WHITE)
    add_textbox(slide, "High-quality starting plan built\nbefore any dynamic event occurs",
                6.12, 5.04, 3.5, 0.55, font_size=11,
                color=RGBColor(0xBB, 0xCC, 0xFF))

    add_textbox(slide,
        "ACO gives us a strong starting point — so online repairs only need to fix small problems.",
        0.4, 5.8, 9.2, 0.45, font_size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def slide_fuzzy(prs):
    """Fuzzy logic — adaptive control policy for repair time allocation."""
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_area(slide, "Technique 2: Fuzzy Logic — Adaptive Control Policy",
                   "Maps disruption severity to repair time ceiling — proportional, interpretable, threshold-free")

    # ── Left column: why + design choices ──────────────────────────────
    add_rect(slide, 0.3, 1.25, 4.4, 0.34, NAVY)
    add_textbox(slide, "Design: 3 rules · triangular MFs · Mamdani",
                0.42, 1.29, 4.2, 0.24, font_size=11, bold=True, color=WHITE)

    add_rect(slide, 0.3, 1.59, 4.4, 2.10, LIGHT_BG, NAVY)
    add_textbox(slide,
        "Why triangular MFs?\n"
        "  · Simplest shape with a clear peak — easy\n"
        "    to interpret and tune by domain experts.\n"
        "  · Standard choice for Mamdani FIS.\n\n"
        "Why only 3 rules?\n"
        "  · Occam's razor: 3 linguistic categories\n"
        "    (LOW / MED / HIGH) cover the severity\n"
        "    range with minimal complexity.\n"
        "  · More rules → over-fitting to one event mix.",
        0.42, 1.68, 4.2, 1.92, font_size=10, color=BLACK)

    # Severity → budget curve (ASCII visual)
    add_rect(slide, 0.3, 3.72, 4.4, 0.30, NAVY)
    add_textbox(slide, "Severity → budget (continuous, not threshold)",
                0.42, 3.75, 4.2, 0.22, font_size=10, bold=True, color=WHITE)
    add_rect(slide, 0.3, 4.02, 4.4, 1.76, LIGHT_BG, NAVY)
    add_textbox(slide,
        "budget\n"
        " 80ms ┤                         ╱‾‾‾\n"
        " 61ms ┤               ╱‾‾╲\n"
        " 35ms ┤  ‾‾‾╲\n"
        "      └──────┼──────┼──────▶  severity\n"
        "            0.3    0.6    1.0\n\n"
        "Weighted-average defuzz → smooth curve\n"
        "(no discontinuity at threshold boundaries)",
        0.42, 4.10, 4.2, 1.58, font_size=9, color=BLACK)

    # ── Right column: rules + examples ─────────────────────────────────
    rules = [
        ("Rule 1 — LOW severity",    "IF sev < 0.3  →  TIGHT   (30–40 ms)",
         "new_customer: insert 1 stop.\nMinimal repair needed.", BLUE),
        ("Rule 2 — MED severity",    "IF sev ≈ 0.5  →  NORMAL  (~50 ms)",
         "traffic_delay ×2: re-route\naffected segments.", GREEN),
        ("Rule 3 — HIGH severity",   "IF sev > 0.7  →  EXTENDED (80 ms)",
         "weather ×1.3: 2-opt all routes.\nMaximum repair budget.", ORANGE),
    ]
    for i, (title, rule, example, color) in enumerate(rules):
        ty = 1.25 + i * 1.56
        add_rect(slide, 4.9, ty, 4.8, 0.32, color)
        add_textbox(slide, title, 5.0, ty + 0.05, 4.6, 0.22,
                    font_size=10, bold=True, color=WHITE)
        add_rect(slide, 4.9, ty + 0.32, 4.8, 0.56, WHITE, color)
        add_textbox(slide, rule, 5.0, ty + 0.38, 4.6, 0.44,
                    font_size=11, bold=True, color=color)
        add_rect(slide, 4.9, ty + 0.88, 4.8, 0.56, LIGHT_BG, color)
        add_textbox(slide, example, 5.0, ty + 0.94, 4.6, 0.44,
                    font_size=10, color=GRAY)

    # ── Bottom bar: honest framing ──────────────────────────────────────
    add_rect(slide, 0.3, 5.82, 9.4, 0.80, NAVY)
    add_textbox(slide,
        "This is an adaptive control policy — not a route optimizer.\n"
        "It decides HOW MUCH time to give the repair algorithm per event.\n"
        "At n=100 the ceiling is non-binding (repair takes 2–4 ms). At larger scale, proportional allocation matters.",
        0.5, 5.87, 9.0, 0.70, font_size=11, color=WHITE, align=PP_ALIGN.CENTER)


def slide_ga(prs):
    """Genetic Algorithm — comparison constructor, not a core FACI-DVRP technique."""
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_area(slide, "GA Constructor: Comparison Method",
                   "Used as an alternative initial constructor to benchmark against ACO")

    add_textbox(slide,
        "Just like biological evolution — start with random solutions, mix the best ones, "
        "mutate slightly, repeat. The fittest routes survive.",
        0.4, 1.25, 9.2, 0.52, font_size=13, color=NAVY, align=PP_ALIGN.CENTER)

    # GA steps as a flow
    ga_steps = [
        ("1. Start",    "30 random\nroute orderings\n(chromosomes)"),
        ("2. Select",   "Pick the best 2\nvia tournament\nselection"),
        ("3. Crossover","Mix them with\nOrder Crossover\n(OX operator)"),
        ("4. Mutate",   "Randomly swap\n2 stops (15%\nchance)"),
        ("5. Keep best","Top route always\nsurvives to\nnext generation"),
    ]
    for i, (step, desc) in enumerate(ga_steps):
        lx = 0.25 + i * 1.9
        ty = 2.0
        add_rect(slide, lx, ty, 1.7, 0.38, BLUE)
        add_textbox(slide, step, lx + 0.08, ty + 0.06, 1.55, 0.26,
                    font_size=11, bold=True, color=WHITE)
        add_rect(slide, lx, ty + 0.38, 1.7, 1.5, LIGHT_BG, BLUE)
        add_textbox(slide, desc, lx + 0.1, ty + 0.5, 1.55, 1.28,
                    font_size=10, color=BLACK)
        if i < 4:
            add_textbox(slide, "→", lx + 1.72, ty + 0.72, 0.22, 0.4,
                        font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_textbox(slide, "Repeat for 50 generations", 0.4, 3.98, 9.2, 0.35,
                font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

    # ACO vs GA comparison
    add_rect(slide, 0.3, 4.45, 9.4, 0.38, NAVY)
    add_textbox(slide, "ACO vs GA  —  How do they compare on C101?",
                0.42, 4.5, 9.0, 0.28, font_size=12, bold=True, color=WHITE)

    compare = [
        ("",            "ACO + Fuzzy + Repair",  "GA + Fuzzy + Repair"),
        ("Final cost",  "1988",                  "2426"),
        ("Avg latency", "2.4 ms",                "2.4 ms"),
        ("Acceptance",  "100%",                  "100%"),
        ("Stability",   "0.988",                 "0.988"),
    ]
    col_x = [0.3, 3.4, 6.7]
    col_w = [3.0, 3.2, 3.0]
    for r, row_data in enumerate(compare):
        ty = 4.9 + r * 0.42
        for c, (cell, x, w) in enumerate(zip(row_data, col_x, col_w)):
            bg = NAVY if r == 0 else (RGBColor(0xE8, 0xF0, 0xFF) if c == 1 else WHITE)
            fc = WHITE if r == 0 else (BLUE if c == 1 else BLACK)
            add_rect(slide, x, ty, w - 0.05, 0.38, bg, RGBColor(0xCC, 0xCC, 0xCC))
            add_textbox(slide, cell, x + 0.08, ty + 0.06, w - 0.2, 0.26,
                        font_size=11, bold=(r == 0), color=fc, align=PP_ALIGN.CENTER)


def slide_dataset(prs):
    """Solomon benchmark dataset explanation."""
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_area(slide, "The Dataset: Solomon Benchmark",
                   "Gold standard for VRP research since 1987 — used in every major routing paper")

    # What is Solomon
    add_textbox(slide,
        "100 customers on a 2D map. Each customer has a location (x, y) and a demand (cargo size). "
        "One depot. Fleet of 25 vehicles, each with capacity 200.",
        0.4, 1.25, 9.2, 0.45, font_size=12, color=NAVY, align=PP_ALIGN.CENTER)

    # Three instances side by side
    instances = [
        (BLUE,   "C101",  "Clustered",
         "Customers grouped\nin geographic pockets\n(like city neighborhoods)",
         "100 customers\nAvg demand: 18.1\nTotal demand: 1810",
         "Best case for ACO —\npheromone captures\ncluster structure well"),
        (GREEN,  "R101",  "Random",
         "Customers uniformly\nscattered across\nthe map",
         "100 customers\nAvg demand: 14.6\nTotal demand: 1458",
         "Hardest for routing —\nno obvious clusters\nto exploit"),
        (ORANGE, "RC101", "Mixed",
         "Some clusters +\nsome random scatter\n(suburban-style)",
         "100 customers\nAvg demand: 17.2\nTotal demand: 1724",
         "Intermediate —\npartial structure\nfor heuristics to use"),
    ]
    for i, (color, name, layout, desc, stats, note) in enumerate(instances):
        lx = 0.3 + i * 3.22
        ty = 1.82
        add_rect(slide, lx, ty, 3.0, 0.38, color)
        add_textbox(slide, f"{name}  —  {layout}", lx + 0.1, ty + 0.06, 2.8, 0.26,
                    font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, lx, ty + 0.38, 3.0, 0.88, LIGHT_BG, color)
        add_textbox(slide, desc, lx + 0.12, ty + 0.46, 2.76, 0.72,
                    font_size=10, color=BLACK, align=PP_ALIGN.CENTER)
        add_rect(slide, lx, ty + 1.26, 3.0, 0.72, WHITE, color)
        add_textbox(slide, stats, lx + 0.12, ty + 1.32, 2.76, 0.60,
                    font_size=10, color=color, align=PP_ALIGN.CENTER)
        add_rect(slide, lx, ty + 1.98, 3.0, 0.60, RGBColor(0xF0, 0xF4, 0xFF), color)
        add_textbox(slide, note, lx + 0.12, ty + 2.04, 2.76, 0.48,
                    font_size=9, color=GRAY, align=PP_ALIGN.CENTER)

    # How the simulation uses the data
    add_rect(slide, 0.3, 4.55, 9.4, 0.34, NAVY)
    add_textbox(slide, "How we use it in the dynamic simulation:",
                0.42, 4.60, 9.0, 0.24, font_size=11, bold=True, color=WHITE)

    sim_pts = [
        ("70 customers", "known at start → ACO builds initial routes offline"),
        ("30 customers", "arrive dynamically one by one → 30 new_customer events"),
        ("+ 2 events",   "1 traffic delay (factor 2.0×)  +  1 weather disruption (factor 1.3×)"),
        ("= 32 events",  "per run · seed=42 · same stream for all 6 methods (fair comparison)"),
    ]
    for i, (label, text) in enumerate(sim_pts):
        ty = 5.0 + i * 0.34
        add_textbox(slide, label, 0.4, ty, 1.4, 0.30,
                    font_size=11, bold=True, color=NAVY)
        add_textbox(slide, text,  1.85, ty, 7.8, 0.30,
                    font_size=11, color=BLACK)

    add_rect(slide, 0.3, 6.44, 9.4, 0.28, LIGHT_BG, NAVY)
    add_textbox(slide,
        "Three representative instances — different customer distributions — let us check whether findings hold across problem structures.",
        0.5, 6.47, 9.0, 0.22, font_size=10, color=NAVY, align=PP_ALIGN.CENTER, bold=True)


def slide_experiment_setup(prs):
    """How we tested it — the benchmark setup."""
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_area(slide, "How We Tested It",
                   "32 events per run · 3 instance types · 6 methods · seed=42")

    # Instance description
    add_textbox(slide,
        "Each Solomon instance (C101 / R101 / RC101) runs the same dynamic simulation: "
        "70 known customers at start, 30 arriving dynamically, plus 2 disruption events.",
        0.4, 1.25, 9.2, 0.52, font_size=12, color=NAVY, align=PP_ALIGN.CENTER)

    # Simulation setup as a timeline-style layout
    timeline = [
        ("Start of day",        "70 customers\nknown upfront",
         "ACO builds the\ninitial route plan.", BLUE),
        ("Event 1 (×30)",       "New customer\narrives",
         "Fuzzy → 35ms ceiling.\nInsertion heuristic.", GREEN),
        ("Event 31",            "Traffic delay\n(factor 2.0×)",
         "Fuzzy → 61ms ceiling.\n2-opt on affected routes.", ORANGE),
        ("Event 32",            "Weather disruption\n(factor 1.3×)",
         "Fuzzy → 80ms ceiling.\n2-opt on ALL routes.", RGBColor(0x1A, 0x6A, 0x8A)),
    ]

    for i, (when, event, action, color) in enumerate(timeline):
        lx = 0.3 + i * 2.38
        ty = 2.0
        add_rect(slide, lx, ty, 2.2, 0.34, color)
        add_textbox(slide, when, lx + 0.08, ty + 0.05, 2.05, 0.24,
                    font_size=10, bold=True, color=WHITE)
        add_rect(slide, lx, ty + 0.34, 2.2, 1.05, LIGHT_BG, color)
        add_textbox(slide, event, lx + 0.1, ty + 0.44, 2.05, 0.55,
                    font_size=11, bold=True, color=color)
        add_rect(slide, lx, ty + 1.39, 2.2, 0.88, RGBColor(0xE0, 0xE8, 0xF8), color)
        add_textbox(slide, action, lx + 0.1, ty + 1.48, 2.05, 0.72,
                    font_size=10, color=BLACK)

    # Summary stats
    stats = [
        ("32", "Total events"),
        ("30", "Dynamic customers"),
        ("2",  "Disruption events"),
        ("25", "Vehicles"),
        ("6",  "Methods compared"),
    ]
    for i, (val, label) in enumerate(stats):
        lx = 0.3 + i * 1.9
        ty = 4.55
        add_rect(slide, lx, ty, 1.72, 0.85, NAVY)
        add_textbox(slide, val,   lx, ty + 0.05, 1.72, 0.48,
                    font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, label, lx, ty + 0.54, 1.72, 0.28,
                    font_size=9, color=RGBColor(0xBB, 0xCC, 0xFF), align=PP_ALIGN.CENTER)

    add_textbox(slide,
        "All 6 methods run on the exact same event stream with the same random seed (42) — "
        "ensuring a perfectly fair comparison.  Solomon instances used without enforcing time windows (capacity constraints only).",
        0.4, 5.56, 9.2, 0.5, font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

    add_textbox(slide, "Seed: 42  ·  Base budget: 50ms  ·  Dynamic ratio: 30%",
                0.4, 6.12, 9.2, 0.35, font_size=11, color=NAVY,
                align=PP_ALIGN.CENTER, bold=True)


def slide_results(prs):
    """Results with multi-instance data — C101, R101, RC101."""
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_area(slide, "The Results",
                   "C101 · R101 · RC101  ·  32 events each  ·  seed=42  ·  6 methods")

    # Actual numbers from experiments/results/multi/summary_table.txt (updated)
    # All ACO methods produce identical costs — budget policy is cost-neutral at n=100
    headers = ["Method", "C101\nCost", "R101\nCost", "RC101\nCost", "Stability"]
    rows = [
        # method,              C101,   R101,   RC101,  stab
        ["FACI-DVRP (ACO+Fuzzy)",    "1988", "1698", "2135", "0.988"],
        ["ACO+Threshold",            "1988", "1698", "2135", "0.988"],
        ["GA+Fuzzy",                 "2426", "2167", "2635", "0.988"],
        ["ACO+Repair (fixed 50ms)",  "1988", "1698", "2135", "0.988"],
        ["Full Re-opt (NN)",         "1666", "1527", "1802", "0.900/0.856/0.809"],
        ["Static (70 cust. only)",   "1156", "1085", "1329", "1.000"],
    ]

    col_widths = [2.6, 1.2, 1.2, 1.2, 3.2]
    col_starts = [0.3]
    for w in col_widths[:-1]:
        col_starts.append(col_starts[-1] + w)

    ty = 1.28
    for j, (h, w, lx) in enumerate(zip(headers, col_widths, col_starts)):
        add_rect(slide, lx, ty, w - 0.04, 0.40, NAVY)
        add_textbox(slide, h, lx + 0.04, ty + 0.04, w - 0.08, 0.34,
                    font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Row backgrounds
    row_bgs = [
        RGBColor(0xE8, 0xF0, 0xFF),  # FACI — proposed method
        RGBColor(0xF4, 0xF4, 0xF4),  # Threshold
        RGBColor(0xFF, 0xF0, 0xF0),  # GA — worse
        RGBColor(0xF4, 0xF4, 0xF4),  # ACO fixed
        WHITE,                        # Re-opt
        RGBColor(0xF8, 0xF8, 0xF8),  # Static
    ]
    # All ACO methods tie — highlight the whole group as green
    aco_cost_cells = {"1988", "1698", "2135"}

    for i, (row, bg) in enumerate(zip(rows, row_bgs)):
        ty = 1.68 + i * 0.36
        is_proposed = (i == 0)
        for j, (cell, w, lx) in enumerate(zip(row, col_widths, col_starts)):
            is_aco_tie = (j in (1, 2, 3) and cell in aco_cost_cells and i <= 3)
            if is_aco_tie:
                cell_bg = RGBColor(0xD8, 0xED, 0xD8)
                cell_fg = GREEN
            else:
                cell_bg = bg
                cell_fg = BLACK
            add_rect(slide, lx, ty, w - 0.04, 0.32, cell_bg, RGBColor(0xCC, 0xCC, 0xCC))
            add_textbox(slide, cell, lx + 0.04, ty + 0.05, w - 0.08, 0.22,
                        font_size=8.5, bold=is_proposed, color=cell_fg,
                        align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)

    notes_y = 1.68 + 6 * 0.36 + 0.06
    add_textbox(slide,
        "Green = all ACO methods tie — budget policy is cost-neutral when starting from the same initial solution.  "
        "Static serves only ~70 customers — lower cost reflects fewer stops, not better routing.",
        0.3, notes_y, 9.2, 0.24, font_size=8, color=GRAY)
    add_textbox(slide,
        "Full Re-opt achieves low cost but disrupts 10–19% of assignments per event (stability 0.856–0.900).  "
        "All repair methods: 100% acceptance, 0.988 stability.",
        0.3, notes_y + 0.24, 9.2, 0.24, font_size=8, color=GRAY)

    # Three findings
    findings = [
        ("ACO beats GA\non cost",
         "ACO+Fuzzy vs GA+Fuzzy:\nC101: 1988 vs 2426 (−18%)\nR101: 1698 vs 2167 (−22%)\nRC101: 2135 vs 2635 (−19%)"),
        ("Budget policy\ncost-neutral",
         "All ACO methods: identical cost.\nFixed = Fuzzy = Threshold.\nBudget allocates time, not quality."),
        ("Stability &\nacceptance robust",
         "All repair methods: 100%\nacceptance, 0.988 stability\nacross all 3 instances."),
    ]
    colors_f = [BLUE, ORANGE, GREEN]
    for i, ((val, desc), col) in enumerate(zip(findings, colors_f)):
        lx = 0.3 + i * 3.22
        ty = notes_y + 0.54
        add_rect(slide, lx, ty, 3.0, 0.56, col)
        add_textbox(slide, val, lx + 0.1, ty + 0.06, 2.8, 0.44,
                    font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, lx, ty + 0.56, 3.0, 1.0, LIGHT_BG, col)
        add_textbox(slide, desc, lx + 0.12, ty + 0.64, 2.76, 0.86,
                    font_size=9, color=BLACK, align=PP_ALIGN.LEFT)

    add_textbox(slide,
        "ACO constructor > GA on cost (18–22%). All ACO budget variants tie — budget policy allocates time, not cost.",
        0.4, notes_y + 1.62, 9.2, 0.30, font_size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def slide_ablation(prs):
    """Ablation study — component-wise contribution."""
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_area(slide, "Ablation Study: What Each Component Contributes",
                   "Same ACO constructor · same event stream · 4 configurations")

    add_textbox(slide,
        "Controlled experiment: each config uses seed=42 before its own ACO call, "
        "so all four start from the same initial random state. Cost differences isolate repair strategy.",
        0.4, 1.25, 9.2, 0.4, font_size=11, color=NAVY, align=PP_ALIGN.CENTER)

    # Table header
    headers = ["Configuration", "C101", "R101", "RC101", "Stability"]
    col_widths = [4.0, 1.3, 1.3, 1.3, 1.3]
    col_starts = [0.3]
    for w in col_widths[:-1]:
        col_starts.append(col_starts[-1] + w)

    ty = 1.72
    for h, w, lx in zip(headers, col_widths, col_starts):
        add_rect(slide, lx, ty, w - 0.04, 0.38, NAVY)
        add_textbox(slide, h, lx + 0.05, ty + 0.06, w - 0.1, 0.26,
                    font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    rows = [
        ("(1) ACO only — static",              "1156.0", "1085.0", "1328.6", "1.000"),
        ("(2) + Insertion, fixed, no 2-opt",   "2138.7", "2005.3", "2449.4", "1.000"),
        ("(3) + Fuzzy budget, no 2-opt",       "2138.7", "2005.3", "2449.4", "1.000"),
        ("(4) + 2-opt + Relocate (FACI-DVRP)*","1988.4", "1697.5", "2135.2", "0.988"),
    ]
    row_bgs = [
        RGBColor(0xF4, 0xF4, 0xF4),
        RGBColor(0xE8, 0xF0, 0xFF),
        RGBColor(0xD8, 0xEA, 0xFF),
        RGBColor(0xD0, 0xED, 0xD0),  # FACI row — green
    ]
    for i, (row, bg) in enumerate(zip(rows, row_bgs)):
        ty_r = 2.10 + i * 0.44
        is_faci = (i == 3)
        for j, (cell, w, lx) in enumerate(zip(row, col_widths, col_starts)):
            add_rect(slide, lx, ty_r, w - 0.04, 0.40, bg,
                     RGBColor(0xCC, 0xCC, 0xCC))
            add_textbox(slide, cell, lx + 0.06, ty_r + 0.07, w - 0.12, 0.26,
                        font_size=10, bold=is_faci,
                        color=GREEN if is_faci else BLACK,
                        align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)

    # Delta annotations
    deltas = [
        ("Config 1→2: Serving 30 new customers raises cost (+920 to +1121). "
         "Insertion alone achieves 100% acceptance.", BLUE),
        ("Config 2→3: Fuzzy budget has ZERO cost effect without 2-opt (Δ=0 on all instances). "
         "Fuzzy is a policy tool, not a quality tool.", ORANGE),
        ("Config 3→4: 2-opt + cross-route relocate delivers ALL the cost improvement "
         "(−150 to −314). This is the sole quality-driving component.", GREEN),
    ]
    for i, (text, color) in enumerate(deltas):
        ty_d = 3.98 + i * 0.60
        add_rect(slide, 0.3, ty_d, 0.18, 0.44, color)
        add_rect(slide, 0.52, ty_d, 9.1, 0.44, LIGHT_BG, color)
        add_textbox(slide, text, 0.65, ty_d + 0.08, 8.85, 0.30,
                    font_size=10, color=BLACK)

    add_rect(slide, 0.3, 5.62, 9.4, 0.22, LIGHT_BG, NAVY)
    add_textbox(slide,
        "* All configs use random.Random(42) as isolated ACO seed — same as main experiment. Ablation numbers match Table I exactly.",
        0.5, 5.65, 9.0, 0.16, font_size=8.5, color=GRAY, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.3, 5.86, 9.4, 0.36, NAVY)
    add_textbox(slide,
        "Takeaway: Local search drives cost quality. Fuzzy budget controls allocation policy. Both serve different purposes.",
        0.5, 5.92, 9.0, 0.26, font_size=11, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)


def slide_robustness(prs):
    """Robustness / sensitivity analysis."""
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_area(slide, "Robustness Analysis: Does It Hold When Assumptions Change?",
                   "Three sensitivity studies on C101 — budget, intensity, fuzzy boundaries")

    studies = [
        (BLUE,
         "B1 — Base Budget Sensitivity",
         "Varied base budget: 20ms → 100ms",
         "Cost: always 1988.4 (Δ=0 fuzzy vs fixed)\n"
         "at every budget level tested.\n"
         "Budget ceiling is non-binding at n=100:\n"
         "repair completes in 2–4ms regardless.\n"
         "Latency differs by <0.5ms between methods.",
         "Feasible always"),
        (GREEN,
         "B2 — Event Intensity",
         "Varied dynamic ratio: 10% → 50%\n(10 to 50 events)",
         "Cost: identical fuzzy vs fixed at every ratio.\n"
         "10%: 1711.4  20%: 1723.7\n"
         "30%: 1988.4  40%: 2043.8  50%: 2097.1\n"
         "Cost grows with ratio. Stability 0.988+,\n"
         "acceptance 100% at every intensity.",
         "Stable under heavy load"),
        (ORANGE,
         "B3 — MF Boundary Sensitivity",
         "Shifted fuzzy MF peaks ±0.10\nfrom designed values",
         "Cost: invariant at 1988.4 across all shifts.\n"
         "Latency: 2.63–2.70ms (minor variation).\n"
         "Stability exactly 0.988 at all shifts.\n"
         "Acceptance always 100%.\n"
         "MF design is not brittle.",
         "Robust to MF tuning"),
    ]

    for i, (color, title, setup, result, verdict) in enumerate(studies):
        lx = 0.3 + i * 3.22
        ty = 1.25
        add_rect(slide, lx, ty, 3.0, 0.36, color)
        add_textbox(slide, title, lx + 0.1, ty + 0.05, 2.8, 0.26,
                    font_size=11, bold=True, color=WHITE)
        add_rect(slide, lx, ty + 0.36, 3.0, 0.60, RGBColor(0xEE, 0xF3, 0xFF), color)
        add_textbox(slide, setup, lx + 0.1, ty + 0.42, 2.8, 0.48,
                    font_size=9.5, color=GRAY)
        add_rect(slide, lx, ty + 0.96, 3.0, 2.4, LIGHT_BG, color)
        add_textbox(slide, result, lx + 0.1, ty + 1.04, 2.8, 2.2,
                    font_size=10, color=BLACK)
        add_rect(slide, lx, ty + 3.36, 3.0, 0.36, color)
        add_textbox(slide, f"✓  {verdict}", lx + 0.1, ty + 3.42, 2.8, 0.26,
                    font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_rect(slide, 0.3, 5.22, 9.4, 0.06, NAVY)
    add_rect(slide, 0.3, 5.32, 9.4, 0.90, NAVY)
    add_textbox(slide,
        "Takeaway: Cost is invariant to budget policy at n=100 — repair always completes within budget.\n"
        "Both fuzzy and fixed achieve identical cost across all budget levels and event intensities.\n"
        "Stability (0.988) and acceptance (100%) are invariant to all parameter changes.",
        0.5, 5.36, 9.0, 0.82, font_size=11, color=WHITE, align=PP_ALIGN.CENTER)


def slide_key_insight(prs):
    """Why fuzzy still matters despite identical cost at n=100."""
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_area(slide, "Why Fuzzy Budget Control Matters",
                   "At n=100 the ceiling is non-binding — that is a feature, not a flaw")

    # Top honest statement
    add_rect(slide, 0.3, 1.25, 9.4, 0.46, NAVY)
    add_textbox(slide,
        "Ablation result: all ACO budget variants produce identical cost on n=100 instances. "
        "2-opt local search is the sole cost-reducing component.",
        0.5, 1.31, 9.0, 0.36, font_size=11.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Left: Why identical cost is NOT a failure
    add_rect(slide, 0.3, 1.84, 4.55, 0.38, GREEN)
    add_textbox(slide, "Why this is the correct result",
                0.42, 1.90, 4.3, 0.26, font_size=12, bold=True, color=WHITE)
    add_rect(slide, 0.3, 2.22, 4.55, 2.20, RGBColor(0xE8, 0xF8, 0xEC), GREEN)
    add_textbox(slide,
        "Repair completes in 2–4 ms per event.\n"
        "Every budget tier (35 / 61 / 80 ms) is\n"
        "far above that 2–4 ms actual work.\n\n"
        "So the ceiling never binds — all tiers\n"
        "give repair the same effective time.\n\n"
        "Identical cost is not a coincidence:\n"
        "it confirms the n=100 instances are\n"
        "too small for the budget to matter yet.",
        0.44, 2.32, 4.28, 2.04, font_size=10.5, color=BLACK)

    # Arrow
    add_textbox(slide, "→", 4.90, 3.15, 0.42, 0.42,
                font_size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Right: Where fuzzy becomes load-bearing
    add_rect(slide, 5.38, 1.84, 4.32, 0.38, BLUE)
    add_textbox(slide, "Where fuzzy becomes load-bearing",
                5.50, 1.90, 4.1, 0.26, font_size=12, bold=True, color=WHITE)
    add_rect(slide, 5.38, 2.22, 4.32, 2.20, LIGHT_BG, BLUE)
    add_textbox(slide,
        "At n=500–1000 customers, repair\n"
        "work scales as O(m²) per route.\n"
        "Actual repair time rises into the\n"
        "tens-of-ms range.\n\n"
        "At that scale, giving 80ms to a\n"
        "minor event wastes time; giving 35ms\n"
        "to severe weather cuts off the search\n"
        "too early.\n\n"
        "Proportional allocation becomes the\n"
        "binding architectural constraint.",
        5.52, 2.32, 4.10, 2.04, font_size=10.5, color=BLACK)

    # Three reasons fuzzy is still the right design
    add_rect(slide, 0.3, 4.56, 9.4, 0.36, NAVY)
    add_textbox(slide, "Why fuzzy is still the right design choice:",
                0.42, 4.62, 9.0, 0.24, font_size=11, bold=True, color=WHITE)

    reasons = [
        (GREEN,  "Proportional\nby construction",
         "Scales continuously with severity.\nNo hard thresholds to tune."),
        (ORANGE, "Interpretable\npolicy",
         "3 readable IF-THEN rules.\nAuditable and tunable."),
        (BLUE,   "Future-proof\narchitecture",
         "Drop in larger instances:\nbenefit appears automatically."),
    ]
    for i, (color, title, desc) in enumerate(reasons):
        lx = 0.3 + i * 3.22
        ty = 5.00
        add_rect(slide, lx, ty, 3.0, 0.58, color)
        add_textbox(slide, title, lx + 0.1, ty + 0.07, 2.82, 0.46,
                    font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, lx, ty + 0.58, 3.0, 0.60, LIGHT_BG, color)
        add_textbox(slide, desc, lx + 0.12, ty + 0.64, 2.78, 0.50,
                    font_size=10, color=BLACK, align=PP_ALIGN.CENTER)

    add_rect(slide, 0.3, 6.30, 9.4, 0.40, LIGHT_BG, NAVY)
    add_textbox(slide,
        "Take-away: fuzzy is a resource allocation policy, not a route optimizer. "
        "At n=100 it is overhead-free. At n=500+ it becomes the constraint that "
        "keeps latency proportional to event severity.",
        0.5, 6.36, 9.0, 0.30, font_size=10, color=NAVY, align=PP_ALIGN.CENTER)


def slide_conclusion(prs):
    """Clean combined conclusion + next steps slide."""
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_area(slide, "Conclusion & Future Work",
                   "FACI-DVRP: three CI techniques, one real problem")

    # ── Three key findings ──────────────────────────────────────────────
    findings = [
        (GREEN,  "Fuzzy logic as\ncontrol policy",
                 "Maps severity → repair time ceiling.\n35ms (simple) to 80ms (severe).\nProportional, interpretable, threshold-free."),
        (BLUE,   "ACO outperforms\nGA on cost",
                 "ACO+Fuzzy beats GA+Fuzzy by\n18–22% across all 3 instances.\nBetter initial routes matter."),
        (ORANGE, "2-opt drives\ncost reduction",
                 "Budget policy is cost-neutral at n=100.\nLocal search is the sole quality driver.\nAblation confirms Δ=0 without 2-opt."),
    ]
    for i, (color, title, desc) in enumerate(findings):
        lx = 0.3 + i * 3.22
        ty = 1.28
        add_rect(slide, lx, ty, 3.0, 0.52, color)
        add_textbox(slide, title, lx + 0.1, ty + 0.08, 2.8, 0.38,
                    font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, lx, ty + 0.52, 3.0, 1.05, LIGHT_BG, color)
        add_textbox(slide, desc, lx + 0.12, ty + 0.60, 2.76, 0.92,
                    font_size=10, color=BLACK)

    # ── Divider ─────────────────────────────────────────────────────────
    add_rect(slide, 0.3, 2.98, 9.4, 0.04, RGBColor(0xCC, 0xCC, 0xCC))

    # ── Future work — 3 compact items ───────────────────────────────────
    add_textbox(slide, "Future Work", 0.3, 3.10, 2.0, 0.30,
                font_size=12, bold=True, color=NAVY)

    future = [
        (BLUE,   "Larger instances (n=500+)",
                 "2-opt scales O(m²) — at that size repair takes tens of ms "
                 "and the fuzzy ceiling becomes binding. "
                 "Proportionality advantage will be empirically measurable."),
        (GREEN,  "All 56 Solomon instances",
                 "Extend benchmarks across full C1/R1/RC1 families. "
                 "Report mean ± std to confirm generalizability of ACO vs GA finding."),
        (ORANGE, "Enforce time windows",
                 "Add VRPTW feasibility checks inside repair loop. "
                 "Budget ceiling will bind sooner, making the fuzzy tier gap significant."),
    ]
    for i, (color, title, desc) in enumerate(future):
        lx = 0.3 + i * 3.22
        ty = 3.46
        add_rect(slide, lx, ty, 3.0, 0.36, color)
        add_textbox(slide, title, lx + 0.1, ty + 0.06, 2.8, 0.26,
                    font_size=11, bold=True, color=WHITE)
        add_rect(slide, lx, ty + 0.36, 3.0, 0.88, LIGHT_BG, color)
        add_textbox(slide, desc, lx + 0.12, ty + 0.44, 2.76, 0.76,
                    font_size=9.5, color=BLACK)

    # ── One-sentence takeaway ────────────────────────────────────────────
    add_rect(slide, 0.3, 4.46, 9.4, 0.56, NAVY)
    add_textbox(slide,
        "Fuzzy logic is the right architecture for proportional resource allocation in real-time DVRP — "
        "and ACO builds better initial routes than GA under the same framework.",
        0.5, 4.54, 9.0, 0.42, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ── Limitations (compact) ────────────────────────────────────────────
    add_textbox(slide, "Limitations", 0.3, 5.14, 1.8, 0.28,
                font_size=11, bold=True, color=ORANGE)
    limits = [
        "Budget ceiling non-binding at n=100 — proportionality not yet empirically demonstrated at this scale",
        "3 Solomon instances only — findings need validation across full benchmark suite",
    ]
    for i, lim in enumerate(limits):
        add_textbox(slide, "· " + lim, 0.3, 5.44 + i * 0.30, 9.4, 0.26,
                    font_size=9.5, color=RGBColor(0x88, 0x44, 0x00))

    # ── References ───────────────────────────────────────────────────────
    add_rect(slide, 0.3, 6.12, 9.4, 0.02, RGBColor(0xCC, 0xCC, 0xCC))
    refs = (
        "[1] Dorigo & Gambardella, IEEE TEC, 1997  ·  "
        "[2] Solomon, Operations Research, 1987  ·  "
        "[3] Zadeh, Information & Control, 1965  ·  "
        "[4] Pillac et al., EJOR, 2013  ·  "
        "[5] Psaraftis et al., Networks, 2016"
    )
    add_textbox(slide, refs, 0.3, 6.18, 9.4, 0.42, font_size=8, color=GRAY)


def slide_future(prs):
    pass  # merged into slide_conclusion


def slide_contribution_vs_existing(prs):
    """Crystal-clear slide: what exists vs what is new."""
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_area(slide, "What Exists vs What We Propose",
                   "The one new idea — and the building blocks borrowed from prior work")

    # LEFT — existing
    add_rect(slide, 0.3, 1.25, 4.5, 0.42, GRAY)
    add_textbox(slide, "EXISTING  (prior work, we use these as tools)",
                0.42, 1.3, 4.3, 0.3, font_size=12, bold=True, color=WHITE)

    existing = [
        ("ACO for VRP",            "Dorigo & Gambardella, 1997",
         "Ants + pheromones build routes.\nWell-established metaheuristic."),
        ("GA for VRP (comparison)", "Vidal et al., Comp. & OR, 2013",
         "Used here as a comparison method only.\nNot part of the proposed approach."),
        ("2-opt local search",      "Lin, 1965 / Applegate et al., 2006",
         "Improve routes by swapping\npairs of edges."),
        ("Incremental DVRP repair", "Pillac et al., EJOR, 2013",
         "Fix only affected routes\ninstead of full rebuild."),
    ]
    for i, (name, ref, desc) in enumerate(existing):
        ty = 1.75 + i * 1.2
        add_rect(slide, 0.3, ty, 4.5, 0.3, RGBColor(0xDD, 0xDD, 0xDD))
        add_textbox(slide, name, 0.42, ty + 0.04, 2.2, 0.22,
                    font_size=11, bold=True, color=BLACK)
        add_textbox(slide, ref,  2.6,  ty + 0.04, 2.1, 0.22,
                    font_size=9,  color=GRAY)
        add_rect(slide, 0.3, ty + 0.3, 4.5, 0.82, LIGHT_BG, RGBColor(0xCC, 0xCC, 0xCC))
        add_textbox(slide, desc, 0.42, ty + 0.36, 4.2, 0.72, font_size=10, color=BLACK)

    # RIGHT — proposed
    add_rect(slide, 5.2, 1.25, 4.5, 0.42, GREEN)
    add_textbox(slide, "PROPOSED  (our novel contribution)",
                5.32, 1.3, 4.3, 0.3, font_size=12, bold=True, color=WHITE)

    add_rect(slide, 5.2, 1.75, 4.5, 0.06, GREEN)
    add_rect(slide, 5.2, 1.81, 4.5, 4.3, RGBColor(0xE4, 0xF6, 0xE8), GREEN)

    add_textbox(slide,
        "Fuzzy Logic Time Budget Controller",
        5.35, 1.92, 4.2, 0.42, font_size=16, bold=True, color=GREEN)

    add_textbox(slide,
        "Assigns proportional ceiling per event type:\n"
        "~35ms for simple events, up to 80ms for\n"
        "severe weather disruptions.",
        5.35, 2.42, 4.2, 0.85, font_size=12, color=BLACK)

    add_textbox(slide,
        "NOT a route optimizer — it allocates repair time, not route quality.",
        5.35, 3.3, 4.2, 0.38, font_size=11, bold=True, color=ORANGE)

    add_rect(slide, 5.35, 3.72, 4.1, 0.04, RGBColor(0xAA, 0xDD, 0xAA))

    add_textbox(slide,
        "Fuzzy-based adaptive repair-budget allocation for DVRP — a policy-driven, "
        "interpretable alternative to fixed or threshold budgeting.",
        5.35, 3.84, 4.2, 0.42, font_size=11, color=BLACK)

    examples = [
        ("new_customer",       "0.25", "35ms", "LOW severity   → less time"),
        ("traffic_delay ×2",   "0.65", "61ms", "MED severity  → moderate time"),
        ("weather ×1.3",       "0.80", "80ms", "HIGH severity → more time"),
    ]
    for i, (ev, sev, budget, label) in enumerate(examples):
        ty = 4.38 + i * 0.56
        add_rect(slide, 5.35, ty, 4.1, 0.5, WHITE, GREEN)
        add_textbox(slide, ev,     5.46, ty + 0.04, 1.5, 0.22, font_size=10, color=GRAY)
        add_textbox(slide, f"sev={sev} → {budget}", 5.46, ty + 0.24, 2.1, 0.22,
                    font_size=11, bold=True, color=GREEN)
        add_textbox(slide, label,  7.6,  ty + 0.14, 1.8, 0.25, font_size=9, color=BLACK)

    # Bottom bar — the one-liner
    add_rect(slide, 0.3, 6.1, 9.4, 0.82, NAVY)
    add_textbox(slide,
        "Prior work applies ACO, GA, and repair operators separately.\n"
        "FACI-DVRP combines them with a fuzzy budget controller that treats"
        " time allocation itself as a CI design decision.",
        0.5, 6.17, 9.0, 0.7, font_size=12, color=WHITE, align=PP_ALIGN.CENTER)


def slide_ga_finding(prs):
    """Empirical CI comparison — ACO beats GA on cost."""
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_area(slide, "CI Constructor Comparison",
                   "ACO+Fuzzy outperforms GA+Fuzzy on cost — across all three instance types")

    # Large header
    add_rect(slide, 0.3, 1.25, 9.4, 0.52, NAVY)
    add_textbox(slide, "ACO+Fuzzy  vs  GA+Fuzzy", 0.4, 1.3, 9.2, 0.42,
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Cost comparison table
    headers_ga = ["Instance", "ACO+Fuzzy Cost", "GA+Fuzzy Cost", "Difference"]
    rows_ga = [
        ["C101  (clustered)",  "1988", "2426", "18% cheaper  (ACO wins)"],
        ["R101  (random)",     "1698", "2167", "22% cheaper  (ACO wins)  ← consistent"],
        ["RC101 (mixed)",      "2135", "2635", "19% cheaper  (ACO wins)"],
    ]

    col_widths_ga = [2.4, 2.0, 2.0, 3.2]
    col_starts_ga = [0.3]
    for w in col_widths_ga[:-1]:
        col_starts_ga.append(col_starts_ga[-1] + w)

    ty = 1.9
    for j, (h, w, lx) in enumerate(zip(headers_ga, col_widths_ga, col_starts_ga)):
        add_rect(slide, lx, ty, w - 0.04, 0.38, NAVY)
        add_textbox(slide, h, lx + 0.05, ty + 0.06, w - 0.1, 0.26,
                    font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    row_bgs_ga = [
        RGBColor(0xD8, 0xED, 0xD8),  # C101 — biggest gap, highlight
        RGBColor(0xF2, 0xF8, 0xF2),
        RGBColor(0xF2, 0xF8, 0xF2),
    ]
    for i, (row, bg) in enumerate(zip(rows_ga, row_bgs_ga)):
        ty = 2.28 + i * 0.46
        for j, (cell, w, lx) in enumerate(zip(row, col_widths_ga, col_starts_ga)):
            fc = GREEN if j == 3 else BLACK
            add_rect(slide, lx, ty, w - 0.04, 0.42, bg, RGBColor(0xCC, 0xCC, 0xCC))
            add_textbox(slide, cell, lx + 0.05, ty + 0.08, w - 0.1, 0.28,
                        font_size=11, bold=(i == 0), color=fc, align=PP_ALIGN.CENTER)

    # Why explanation
    add_rect(slide, 0.3, 3.72, 9.4, 0.38, NAVY)
    add_textbox(slide, "Why does ACO win on cost?", 0.42, 3.78, 9.0, 0.26,
                font_size=12, bold=True, color=WHITE)

    add_rect(slide, 0.3, 4.1, 9.4, 1.38, LIGHT_BG, NAVY)
    add_textbox(slide,
        "ACO runs 10 ants x 20 iterations = 200 construction passes, with pheromone reinforcement "
        "guiding each ant toward high-quality arcs discovered by previous ants. This spatial memory "
        "accumulation is highly effective for VRPTW instances.\n\n"
        "GA's greedy capacity-split decoder assigns customers to routes in permutation order — "
        "which may not align with geographic clusters. OX crossover preserves relative ordering "
        "within the permutation, but the split-decoder interpretation limits meaningful exploitation.",
        0.45, 4.18, 9.05, 1.22, font_size=11, color=BLACK)

    # Bottom recommendation
    add_rect(slide, 0.3, 5.58, 9.4, 0.78, NAVY)
    add_textbox(slide,
        "ACO constructor outperforms GA on cost across all instances.\n"
        "Fuzzy budget provides interpretable allocation policy — not a guaranteed cost advantage over fixed or threshold budgets.",
        0.5, 5.68, 9.0, 0.58, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide_multi_fuzzy(prs):
    """Multi-input vs single-input fuzzy controller comparison."""
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_area(slide, "Extended: 3-Input Fuzzy Controller",
                   "Context-aware budget allocation using route state — beyond event severity alone")

    # Header
    add_rect(slide, 0.3, 1.25, 9.4, 0.44, NAVY)
    add_textbox(slide, "Single-Input FIS  vs  3-Input FIS (ACO+Multi-Fuzzy)",
                0.4, 1.3, 9.2, 0.34, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Left: inputs diagram
    add_rect(slide, 0.3, 1.82, 4.4, 0.32, RGBColor(0x26, 0x5E, 0xAD))
    add_textbox(slide, "Single-Input FIS — 1 variable, 3 rules",
                0.4, 1.88, 4.2, 0.22, font_size=10, bold=True, color=WHITE)

    inputs_single = [
        "Input:   disruption_severity  (0–1)",
        "Rules:   LOW→TIGHT  MED→NORMAL  HIGH→EXTENDED",
        "Output:  time_budget_ms",
    ]
    for k, txt in enumerate(inputs_single):
        add_textbox(slide, txt, 0.38, 2.17 + k * 0.32, 4.2, 0.28, font_size=10, color=BLACK)

    add_rect(slide, 0.3, 3.14, 4.4, 0.32, RGBColor(0x1A, 0x7A, 0x3C))
    add_textbox(slide, "3-Input FIS — 3 variables, 12 rules",
                0.4, 3.2, 4.2, 0.22, font_size=10, bold=True, color=WHITE)

    inputs_multi = [
        "Input 1: disruption_severity       (0–1)",
        "Input 2: affected_route_ratio      (0–1)",
        "Input 3: route_load_factor         (0–1)",
        "Rules:   12 (min-AND Mamdani)",
        "Output:  time_budget_ms",
    ]
    for k, txt in enumerate(inputs_multi):
        add_textbox(slide, txt, 0.38, 3.49 + k * 0.28, 4.2, 0.24, font_size=9.5, color=BLACK)

    # Right: results table
    add_rect(slide, 4.9, 1.82, 4.8, 0.32, NAVY)
    add_textbox(slide, "Cost Results (lower is better)",
                5.0, 1.88, 4.6, 0.22, font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    tbl_headers = ["Instance", "Single-FIS", "3-Input FIS", "Δ"]
    tbl_rows = [
        ["C101  (clustered)", "1988", "1988", "Δ=0 (tie)"],
        ["R101  (random)",    "1698", "1698", "Δ=0 (tie)"],
        ["RC101 (mixed)",     "2135", "2135", "Δ=0 (tie)"],
    ]
    tbl_widths = [1.6, 1.0, 1.1, 1.0]
    tbl_starts = [4.9]
    for w in tbl_widths[:-1]:
        tbl_starts.append(tbl_starts[-1] + w)

    ty = 2.18
    for j, (h, w, lx) in enumerate(zip(tbl_headers, tbl_widths, tbl_starts)):
        add_rect(slide, lx, ty, w - 0.04, 0.32, NAVY)
        add_textbox(slide, h, lx + 0.04, ty + 0.06, w - 0.08, 0.22,
                    font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    row_bgs_m = [
        RGBColor(0xD8, 0xED, 0xD8),
        RGBColor(0xD8, 0xED, 0xD8),
        RGBColor(0xD8, 0xED, 0xD8),
    ]
    for i, (row, bg) in enumerate(zip(tbl_rows, row_bgs_m)):
        ty = 2.54 + i * 0.38
        for j, (cell, w, lx) in enumerate(zip(row, tbl_widths, tbl_starts)):
            fc = GREEN if j == 3 else BLACK
            add_rect(slide, lx, ty, w - 0.04, 0.34, bg, RGBColor(0xCC, 0xCC, 0xCC))
            add_textbox(slide, cell, lx + 0.04, ty + 0.07, w - 0.08, 0.22,
                        font_size=9, bold=(j == 3), color=fc, align=PP_ALIGN.CENTER)

    # Interpretation box
    add_rect(slide, 4.9, 3.68, 4.8, 1.72, LIGHT_BG, NAVY)
    add_textbox(slide,
        "Why do both FIS variants tie?\n\n"
        "Both start from the same ACO solution.\n"
        "At n=100, all repair completes in 2–4ms\n"
        "regardless of the budget ceiling.\n\n"
        "3-input FIS still assigns different ceilings:\n"
        "→ new_customer: tighter (FEW routes)\n"
        "→ weather + heavy: extended (MANY routes)\n\n"
        "Value: richer policy semantics, same cost.",
        5.05, 3.78, 4.55, 1.54, font_size=9.5, color=BLACK)

    # Bottom stability note
    add_rect(slide, 0.3, 5.55, 9.4, 0.52, NAVY)
    add_textbox(slide,
        "Both FIS variants: 100% acceptance · 0.988 stability · 2–4 ms avg latency (wall-clock, averaged over 32 events)\n"
        "No measurable latency difference between variants — both complete well within the allocated budget at n=100",
        0.5, 5.63, 9.0, 0.38, font_size=11, color=WHITE, align=PP_ALIGN.CENTER)


def slide_local_search(prs):
    """CI Technique 3: 2-opt local search + cross-route relocate."""
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_area(slide, "CI Technique 3: 2-opt Local Search + Repair",
                   "The component that drives ALL cost improvement — proven by ablation")

    # Big verdict from ablation
    add_rect(slide, 0.3, 1.25, 9.4, 0.50, ORANGE)
    add_textbox(slide,
        "Ablation result: adding fuzzy budget alone → Δ cost = 0.  Adding 2-opt → 7–14% cost reduction.",
        0.5, 1.33, 9.0, 0.36, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Left: how 2-opt works
    add_rect(slide, 0.3, 1.88, 4.55, 0.36, ORANGE)
    add_textbox(slide, "How 2-opt works",
                0.42, 1.93, 4.3, 0.26, font_size=12, bold=True, color=WHITE)

    add_rect(slide, 0.3, 2.24, 4.55, 2.70, LIGHT_BG, ORANGE)
    add_textbox(slide,
        "Given a route  [A → B → C → D → E → depot]\n\n"
        "Try reversing every sub-segment:\n"
        "  [A → C → B → D → E]  — shorter?\n"
        "  [A → B → D → C → E]  — shorter?\n"
        "  ...all (i, k) pairs...\n\n"
        "Accept the best swap that reduces\n"
        "total route distance.\n\n"
        "Repeat until no improving swap exists\n"
        "OR the fuzzy time budget expires.",
        0.44, 2.34, 4.3, 2.50, font_size=10.5, color=BLACK)

    # Middle arrow
    add_textbox(slide, "→", 4.95, 3.4, 0.4, 0.5,
                font_size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Right: cross-route relocate
    add_rect(slide, 5.4, 1.88, 4.3, 0.36, BLUE)
    add_textbox(slide, "Cross-route relocate (bonus operator)",
                5.52, 1.93, 4.1, 0.26, font_size=12, bold=True, color=WHITE)

    add_rect(slide, 5.4, 2.24, 4.3, 2.70, LIGHT_BG, BLUE)
    add_textbox(slide,
        "After 2-opt, try moving one customer\n"
        "from their current route to a better\n"
        "position in a different route:\n\n"
        "  Route 1: [A → B → C]  remove B\n"
        "  Route 2: [D → E]  insert B → [D → B → E]\n\n"
        "Accept if total cost decreases AND\n"
        "capacity constraints still satisfied.\n\n"
        "Causes stability = 0.988 (not 1.000) —\n"
        "the occasional inter-vehicle swap.",
        5.54, 2.34, 4.1, 2.50, font_size=10.5, color=BLACK)

    # Bottom: when it runs
    add_rect(slide, 0.3, 5.08, 9.4, 0.36, NAVY)
    add_textbox(slide, "When does repair run?",
                0.42, 5.13, 9.0, 0.24, font_size=12, bold=True, color=WHITE)

    triggers = [
        ("new_customer\narrives",
         "Insert at best position\n+ one cross-route relocate\nwithin fuzzy budget"),
        ("traffic_delay\non a segment",
         "2-opt on affected routes\nonly — within fuzzy budget\n(targeted repair)"),
        ("weather_disruption\n(global)",
         "2-opt on ALL non-empty\nroutes + cross-route relocate\nwithin fuzzy budget"),
    ]
    for i, (event, action) in enumerate(triggers):
        lx = 0.3 + i * 3.22
        add_rect(slide, lx, 5.48, 3.0, 0.46, RGBColor(0xEE, 0xF3, 0xFF), NAVY)
        add_textbox(slide, event, lx + 0.1, 5.52, 2.8, 0.38,
                    font_size=9.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_rect(slide, lx, 5.94, 3.0, 0.72, LIGHT_BG, NAVY)
        add_textbox(slide, action, lx + 0.1, 5.98, 2.8, 0.64,
                    font_size=9.5, color=BLACK, align=PP_ALIGN.CENTER)


def slide_architecture(prs):
    """System architecture: Event → Severity → FIS → Budget → Repair → Route."""
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_area(slide, "System Architecture: FACI-DVRP Pipeline",
                   "How a real-time disruption flows through all three CI components")

    # Pipeline boxes: 5 stages
    stages = [
        (ORANGE,                         "Disruption\nEvent",
         "new_customer\ntraffic_delay\nweather_disruption"),
        (RGBColor(0x1A, 0x6A, 0x8A),    "Severity\nScore",
         "Compute event\nseverity [0–1]\nfrom event type"),
        (GREEN,                          "Mamdani\nFIS",
         "3 rules · triangular MFs\nOutputs time ceiling\n(35 / 61 / 80 ms)"),
        (BLUE,                           "2-opt +\nRelocate",
         "Repair affected\nroutes within\nthe budget window"),
        (NAVY,                           "Updated\nRoute Plan",
         "100% acceptance\n0.988 stability\n2–4 ms latency"),
    ]

    n = len(stages)
    box_w = 1.65
    gap = 0.18
    total_w = n * box_w + (n - 1) * gap
    start_x = (10.0 - total_w) / 2.0
    box_top = 1.80
    header_h = 0.58
    body_h = 0.85

    for i, (color, title, detail) in enumerate(stages):
        lx = start_x + i * (box_w + gap)
        add_rect(slide, lx, box_top, box_w, header_h, color)
        add_textbox(slide, title, lx + 0.06, box_top + 0.07, box_w - 0.12, header_h - 0.1,
                    font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, lx, box_top + header_h, box_w, body_h, LIGHT_BG, color)
        add_textbox(slide, detail, lx + 0.07, box_top + header_h + 0.06,
                    box_w - 0.14, body_h - 0.1,
                    font_size=9, color=BLACK, align=PP_ALIGN.CENTER)

        # Arrow between boxes
        if i < n - 1:
            ax = lx + box_w + 0.02
            ay = box_top + (header_h + body_h) / 2 - 0.12
            add_textbox(slide, "→", ax, ay, gap + 0.04, 0.30,
                        font_size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # CI technique labels below each stage
    labels = ["", "CI #1\nACO output", "CI #2\nFuzzy Logic", "CI #3\n2-opt", ""]
    ci_colors = [WHITE, BLUE, GREEN, ORANGE, WHITE]
    for i, (lbl, lc) in enumerate(zip(labels, ci_colors)):
        if not lbl:
            continue
        lx = start_x + i * (box_w + gap)
        add_rect(slide, lx, box_top + header_h + body_h + 0.12, box_w, 0.42, lc)
        add_textbox(slide, lbl, lx + 0.06, box_top + header_h + body_h + 0.17,
                    box_w - 0.12, 0.34,
                    font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Complexity bar
    add_rect(slide, 0.3, 4.20, 9.4, 0.36, NAVY)
    add_textbox(slide, "Computational complexity per event:",
                0.42, 4.26, 2.5, 0.24, font_size=10, bold=True, color=WHITE)
    add_textbox(slide,
        "ACO (offline) — O(n² · n_ants · n_iter)   ·   "
        "Fuzzy FIS — O(1) lookup   ·   "
        "Best-position insertion — O(K·n)   ·   "
        "2-opt per route — O(m²) where m = route length",
        2.95, 4.26, 6.65, 0.24, font_size=9.5, color=RGBColor(0xCC, 0xDD, 0xFF))

    # Key property boxes
    props = [
        (GREEN,  "Event-driven",
         "Repair only triggered\nby actual disruptions\n— zero idle overhead"),
        (BLUE,   "Warm-start reuse",
         "Existing routes are\nkept; only affected\nportions are repaired"),
        (ORANGE, "Budget-bounded",
         "Fuzzy ceiling ensures\nworst-case latency\nremains predictable"),
        (NAVY,   "Proportional",
         "Simple events get\ntight budgets; severe\ndisruptions get more"),
    ]
    for i, (color, title, desc) in enumerate(props):
        lx = 0.3 + i * 2.4
        ty = 4.72
        add_rect(slide, lx, ty, 2.22, 0.36, color)
        add_textbox(slide, title, lx + 0.08, ty + 0.06, 2.08, 0.24,
                    font_size=10, bold=True, color=WHITE)
        add_rect(slide, lx, ty + 0.36, 2.22, 0.90, LIGHT_BG, color)
        add_textbox(slide, desc, lx + 0.1, ty + 0.42, 2.04, 0.78,
                    font_size=9.5, color=BLACK)

    # Bottom note
    add_rect(slide, 0.3, 6.12, 9.4, 0.52, LIGHT_BG, NAVY)
    add_textbox(slide,
        "ACO runs ONCE offline before the day starts — all real-time work is "
        "Severity→FIS→Repair, completing in 2–4 ms per event on n=100 instances.",
        0.5, 6.20, 9.0, 0.38, font_size=11, color=NAVY, align=PP_ALIGN.CENTER)


# -----------------------------------------------------------------------
# Build presentation
# -----------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title(prs)                    # 1  — Title
    slide_hook(prs)                     # 2  — Picture this...
    slide_two_bad_options(prs)          # 3  — Why naive answers fail
    slide_solution_overview(prs)        # 4  — Meet FACI-DVRP
    slide_contribution_vs_existing(prs) # 5  — What's new vs what exists
    slide_architecture(prs)             # 6  — System architecture pipeline
    slide_aco(prs)                      # 7  — Technique 1: ACO
    slide_fuzzy(prs)                    # 8  — Technique 2: Fuzzy Logic
    slide_local_search(prs)             # 9  — Technique 3: 2-opt local search
    slide_ga(prs)                       # 10 — GA constructor (comparison)
    slide_ga_finding(prs)               # 11 — ACO vs GA finding
    slide_dataset(prs)                  # 12 — Solomon benchmark dataset
    slide_experiment_setup(prs)         # 13 — How we tested it
    slide_results(prs)                  # 14 — The results
    slide_ablation(prs)                 # 15 — Ablation study
    slide_key_insight(prs)              # 16 — Why fuzzy still matters
    slide_conclusion(prs)               # 17 — Conclusion + future work (combined)

    out = "presentation.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    build()
